#!/usr/bin/env python3
"""Hybrid Workforce — UI Mockup"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG      = RGBColor(0x0d, 0x0f, 0x1a)
SIDEBAR = RGBColor(0x12, 0x14, 0x22)
CARD    = RGBColor(0x1a, 0x1d, 0x2e)
CARD_LT = RGBColor(0x22, 0x25, 0x3a)
ACCENT  = RGBColor(0x6c, 0x9e, 0xff)
ACCENT2 = RGBColor(0x7c, 0xdb, 0xc4)
WARN    = RGBColor(0xff, 0xb8, 0x6c)
GREEN   = RGBColor(0x6b, 0xe6, 0x9c)
RED     = RGBColor(0xff, 0x7b, 0x7b)
WHITE   = RGBColor(0xf0, 0xf0, 0xf5)
GRAY    = RGBColor(0x8a, 0x8d, 0xa0)
DIM     = RGBColor(0x44, 0x47, 0x5c)
LIGHT   = RGBColor(0xcc, 0xcc, 0xd5)
PURPLE  = RGBColor(0xb0, 0x8e, 0xff)

def bg(slide, color=BG):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def txt(slide, l, t, w, h, text, size=14, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = "Helvetica Neue"; p.alignment = align
    return tf

def rect(slide, l, t, w, h, fill, radius=False, border_color=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    if radius:
        s.adjustments[0] = 0.04
    return s

def dot(slide, l, t, r, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(l), Inches(t), Inches(r), Inches(r))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def divider(slide, l, t, w, color=DIM):
    rect(slide, l, t, w, 0.02, color)

def sidebar(slide):
    """Draw common left sidebar"""
    rect(slide, 0, 0, 1.6, 7.5, SIDEBAR)
    txt(slide, 0.15, 0.25, 1.3, 0.4, "⬡  HW", size=16, color=ACCENT, bold=True)
    divider(slide, 0.15, 0.8, 1.3)
    nav = [("🏠", "Dashboard"), ("👥", "Team"), ("📋", "Tasks"),
           ("💬", "Feed"), ("⚙️", "Settings")]
    for i, (icon, label) in enumerate(nav):
        y = 1.05 + i * 0.7
        txt(slide, 0.2, y, 1.2, 0.4, f"{icon}  {label}", size=13, color=GRAY)
    # active highlight on first item
    rect(slide, 0, 1.0, 1.6, 0.42, RGBColor(0x6c, 0x9e, 0xff) if False else CARD_LT)

def topbar(slide, title, subtitle=""):
    rect(slide, 1.6, 0, 11.733, 0.75, CARD)
    txt(slide, 1.85, 0.1, 8, 0.35, title, size=18, color=WHITE, bold=True)
    if subtitle:
        txt(slide, 1.85, 0.42, 8, 0.28, subtitle, size=12, color=GRAY)
    # Avatar placeholder
    dot(slide, 12.6, 0.18, 0.4, DIM)
    txt(slide, 12.62, 0.22, 0.4, 0.3, "J", size=14, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER)

def slide_label(slide, text):
    rect(slide, 0, 7.2, 13.333, 0.3, RGBColor(0x08, 0x09, 0x12))
    txt(slide, 0.2, 7.22, 13, 0.25, text, size=10, color=DIM)

# ══════════════════════════════════════════════════════════════
# SLIDE 1: Home Dashboard
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide)
sidebar(slide)
topbar(slide, "Dashboard", "你好，Juanjuan — 今天有 2 件事需要你决策")
slide_label(slide, "Screen 1 of 4 — Home Dashboard")

# ── Needs Attention banner ──
rect(slide, 1.7, 0.85, 11.5, 0.7, RGBColor(0x2a, 0x22, 0x10), border_color=WARN)
dot(slide, 1.9, 1.05, 0.25, WARN)
txt(slide, 2.25, 0.93, 9, 0.2, "需要你决策", size=11, color=WARN, bold=True)
txt(slide, 2.25, 1.15, 9, 0.28,
    "pm-Octopus 请求确认 PRD v2 发布  ·  SuperBoss 等待产品方向拍板",
    size=12, color=LIGHT)
txt(slide, 12.5, 1.0, 0.6, 0.4, "查看", size=12, color=ACCENT, bold=True)

# ── Team Status ──
txt(slide, 1.85, 1.75, 4, 0.35, "团队状态", size=14, color=WHITE, bold=True)

agents = [
    ("🐙 pm-Octopus", "正在撰写 PRD v2", GREEN, "进行中"),
    ("🤖 SuperBoss",  "等待方向确认", WARN,  "等待"),
    ("⚡ SuperCrew",  "Milestone 4 开发中", GREEN, "进行中"),
    ("📡 EM Agent",   "下次检查 in 2h", GRAY, "空闲"),
]
for i, (name, status, clr, badge) in enumerate(agents):
    y = 2.2 + i * 0.95
    rect(slide, 1.7, y, 5.2, 0.8, CARD, radius=True)
    dot(slide, 1.95, y + 0.28, 0.22, clr)
    txt(slide, 2.3, y + 0.1, 3.0, 0.3, name, size=13, color=WHITE, bold=True)
    txt(slide, 2.3, y + 0.42, 3.0, 0.28, status, size=12, color=GRAY)
    rect(slide, 5.5, y + 0.22, 1.2, 0.35, CARD_LT, radius=True)
    badge_clr = GREEN if badge == "进行中" else WARN if badge == "等待" else GRAY
    txt(slide, 5.55, y + 0.27, 1.1, 0.28, badge, size=11, color=badge_clr,
        bold=True, align=PP_ALIGN.CENTER)

# ── Recent Activity ──
txt(slide, 7.2, 1.75, 5.8, 0.35, "最近活动", size=14, color=WHITE, bold=True)

activities = [
    (GREEN,  "SuperCrew 完成", "API 接口联调，已推送"),
    (ACCENT, "pm-Octopus 更新", "竞品分析报告 v2 已生成"),
    (WARN,   "验收未通过", "POI 数据覆盖率不足，已上报"),
    (GRAY,   "EM Agent 摘要", "今日进展：3 个任务完成，1 个待决策"),
]
for i, (clr, actor, action) in enumerate(activities):
    y = 2.2 + i * 0.95
    rect(slide, 7.2, y, 5.7, 0.8, CARD, radius=True)
    dot(slide, 7.4, y + 0.28, 0.14, clr)
    txt(slide, 7.65, y + 0.1, 2.0, 0.3, actor, size=12, color=clr, bold=True)
    txt(slide, 7.65, y + 0.4, 4.9, 0.3, action, size=12, color=GRAY)

# ── Quick Stats ──
stats = [("3", "任务进行中", ACCENT), ("2", "等待决策", WARN),
         ("12", "今日完成", GREEN), ("0", "异常告警", GRAY)]
for i, (n, l, c) in enumerate(stats):
    x = 1.7 + i * 2.9
    rect(slide, x, 6.05, 2.7, 1.1, CARD, radius=True)
    txt(slide, x+0.25, 6.15, 2.2, 0.5, n, size=32, color=c, bold=True)
    txt(slide, x+0.25, 6.65, 2.2, 0.3, l, size=12, color=GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 2: Team Builder
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide)
sidebar(slide)
topbar(slide, "Team Builder", "配置你的 Agent 团队成员")
slide_label(slide, "Screen 2 of 4 — Team Builder")

# ── Left: Role Library ──
txt(slide, 1.85, 0.9, 4, 0.35, "角色模板库", size=14, color=WHITE, bold=True)
txt(slide, 1.85, 1.28, 4, 0.28, "选择角色添加到你的团队", size=12, color=GRAY)

roles = [
    ("🐙", "PM Agent",       ACCENT,  "需求分析 · PRD · 竞品"),
    ("🎯", "Research Agent", ACCENT2, "信息收集 · 调研 · 数据"),
    ("⚡", "Dev Agent",      PURPLE,  "编码 · 测试 · 部署"),
    ("🔍", "QA Agent",       WARN,    "验收 · 审查 · 异常发现"),
    ("📡", "EM Agent",       GREEN,   "进度 · 协调 · 摘要推送"),
]
for i, (icon, name, clr, desc) in enumerate(roles):
    y = 1.7 + i * 0.95
    rect(slide, 1.7, y, 4.7, 0.8, CARD, radius=True)
    txt(slide, 1.95, y + 0.2, 0.4, 0.4, icon, size=20, align=PP_ALIGN.CENTER)
    txt(slide, 2.5, y + 0.1, 2.2, 0.3, name, size=13, color=WHITE, bold=True)
    txt(slide, 2.5, y + 0.42, 2.5, 0.28, desc, size=11, color=GRAY)
    # Add button
    rect(slide, 5.8, y + 0.2, 0.5, 0.38, ACCENT, radius=True)
    txt(slide, 5.83, y + 0.25, 0.44, 0.28, "+", size=18, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER)

# ── Right: Current Team Config ──
txt(slide, 7.2, 0.9, 5.8, 0.35, "当前团队配置", size=14, color=WHITE, bold=True)
divider(slide, 7.2, 1.32, 5.7)

current = [
    ("🐙 pm-Octopus", ACCENT,  "行为准则 · 记忆 · 方法论", "已配置", GREEN),
    ("⚡ SuperCrew",  PURPLE,  "编码 · 测试 · 部署工具",    "已配置", GREEN),
    ("📡 EM Agent",   GREEN,   "每 2h 检查 · Cron 定时",    "运行中", ACCENT),
]
for i, (name, clr, skills, status, sclr) in enumerate(current):
    y = 1.55 + i * 1.3
    rect(slide, 7.2, y, 5.7, 1.1, CARD, radius=True)
    txt(slide, 7.45, y + 0.12, 3.5, 0.35, name, size=14, color=clr, bold=True)
    txt(slide, 7.45, y + 0.5, 3.5, 0.28, skills, size=11, color=GRAY)
    rect(slide, 11.5, y + 0.35, 1.1, 0.35, CARD_LT, radius=True)
    txt(slide, 11.52, y + 0.4, 1.06, 0.28, status, size=11, color=sclr,
        bold=True, align=PP_ALIGN.CENTER)

# ── Tools panel ──
txt(slide, 7.2, 5.7, 5.8, 0.35, "已安装工具", size=14, color=WHITE, bold=True)
tools = ["browser-use", "小红书爬虫", "PPT 生成", "高德 API", "+ 安装更多"]
for i, tool in enumerate(tools):
    clr = ACCENT if i < 4 else DIM
    bg_clr = CARD_LT if i < 4 else CARD
    rect(slide, 7.2 + i * 1.22, 6.1, 1.1, 0.45, bg_clr, radius=True)
    txt(slide, 7.22 + i * 1.22, 6.18, 1.06, 0.3, tool, size=10, color=clr,
        align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 3: Agent Profile
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide)
sidebar(slide)
topbar(slide, "pm-Octopus — Agent Profile", "你的 AI 产品经理")
slide_label(slide, "Screen 3 of 4 — Agent Profile")

# ── Left column: Identity ──
rect(slide, 1.7, 0.85, 3.8, 6.45, CARD, radius=True)
txt(slide, 2.2, 1.1, 2.8, 0.8, "🐙", size=48, align=PP_ALIGN.CENTER)
txt(slide, 1.9, 1.9, 3.2, 0.45, "pm-Octopus", size=20, color=WHITE,
    bold=True, align=PP_ALIGN.CENTER)
txt(slide, 1.9, 2.35, 3.2, 0.3, "AI 产品经理", size=13, color=GRAY,
    align=PP_ALIGN.CENTER)
dot(slide, 3.15, 2.4, 0.18, GREEN)

divider(slide, 2.0, 2.8, 2.9)

txt(slide, 2.0, 2.95, 2.9, 0.28, "状态", size=11, color=GRAY)
txt(slide, 2.0, 3.23, 2.9, 0.28, "正在撰写 PRD v2", size=13, color=GREEN, bold=True)

txt(slide, 2.0, 3.6, 2.9, 0.28, "已掌握方法论", size=11, color=GRAY)
methods = ["用户洞察框架", "竞品分析框架", "Prompt工程", "LLM评估", "Agent评估"]
for i, m in enumerate(methods):
    y = 3.88 + i * 0.38
    rect(slide, 2.0, y, 2.9, 0.3, CARD_LT, radius=True)
    txt(slide, 2.1, y + 0.04, 2.7, 0.24, m, size=11, color=ACCENT2)

txt(slide, 2.0, 5.85, 2.9, 0.28, "已安装工具", size=11, color=GRAY)
for i, tool in enumerate(["browser-use", "小红书爬虫", "PPT生成"]):
    rect(slide, 2.0 + i * 1.0, 6.13, 0.9, 0.3, CARD_LT, radius=True)
    txt(slide, 2.02 + i * 1.0, 6.18, 0.86, 0.24, tool, size=9, color=ACCENT)

# ── Middle: Memory & Rules ──
txt(slide, 5.8, 0.9, 3.7, 0.35, "行为准则（节选）", size=14, color=WHITE, bold=True)
rect(slide, 5.8, 1.3, 3.7, 2.8, CARD, radius=True)
rules = [
    "✦  任务创建专属文件夹，过程留文档",
    "✦  决策给选项+分析，不甩开放问题",
    "✦  角色边界：不管代码不看PR",
    "✦  需要决策时 @Juanjuan，附背景",
    "✦  方法论遇到新场景自主更新",
]
for i, r in enumerate(rules):
    txt(slide, 6.0, 1.45 + i * 0.48, 3.3, 0.38, r, size=12, color=LIGHT)

txt(slide, 5.8, 4.3, 3.7, 0.35, "近期记忆", size=14, color=WHITE, bold=True)
mems = [
    (ACCENT, "今天", "PRD v2 撰写中，等待 Juanjuan 确认"),
    (GRAY,   "昨天", "完成竞品分析，发现3个差异化机会"),
    (GRAY,   "4/10", "与 SuperBoss 对齐 Milestone 5 计划"),
]
for i, (clr, day, content) in enumerate(mems):
    y = 4.7 + i * 0.65
    rect(slide, 5.8, y, 3.7, 0.55, CARD, radius=True)
    dot(slide, 5.98, y + 0.18, 0.16, clr)
    txt(slide, 6.22, y + 0.06, 0.6, 0.22, day, size=10, color=clr, bold=True)
    txt(slide, 6.22, y + 0.28, 3.1, 0.22, content, size=11, color=GRAY)

# ── Right: Growth & Feedback ──
txt(slide, 9.8, 0.9, 3.3, 0.35, "成长记录", size=14, color=WHITE, bold=True)
rect(slide, 9.8, 1.3, 3.3, 2.3, CARD, radius=True)
growths = [
    (GREEN, "自主优化", "腰部帖子策略，提升需求挖掘质量"),
    (ACCENT, "边界内化", "被纠正后自己写进行为准则"),
    (PURPLE, "方法论", "主动搜索沉淀 5 套方法论"),
]
for i, (clr, title, desc) in enumerate(growths):
    y = 1.45 + i * 0.65
    dot(slide, 9.98, y + 0.14, 0.14, clr)
    txt(slide, 10.2, y + 0.04, 0.9, 0.24, title, size=11, color=clr, bold=True)
    txt(slide, 10.2, y + 0.28, 2.7, 0.24, desc, size=11, color=GRAY)

txt(slide, 9.8, 3.8, 3.3, 0.35, "给 Feedback", size=14, color=WHITE, bold=True)
rect(slide, 9.8, 4.2, 3.3, 1.3, CARD_LT, radius=True,
     border_color=ACCENT)
txt(slide, 10.0, 4.35, 2.9, 0.6,
    "输入反馈，Agent 自动更新行为准则...",
    size=12, color=DIM)
rect(slide, 9.8, 5.55, 3.3, 0.45, ACCENT, radius=True)
txt(slide, 9.8, 5.6, 3.3, 0.35, "发送 → 自动沉淀", size=13, color=WHITE,
    bold=True, align=PP_ALIGN.CENTER)

txt(slide, 9.8, 6.15, 3.3, 0.28, "最近 30 天：收到 12 条 feedback，0 次重犯",
    size=11, color=GREEN)


# ══════════════════════════════════════════════════════════════
# SLIDE 4: Task Feed & Decision Center
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide)
sidebar(slide)
topbar(slide, "Feed", "实时协作动态")
slide_label(slide, "Screen 4 of 4 — Task Feed & Decision Center")

# ── Left: Feed ──
txt(slide, 1.85, 0.9, 5.5, 0.35, "协作动态", size=14, color=WHITE, bold=True)

feeds = [
    (WARN,   "需要决策",  "pm-Octopus",
     "PRD v2 已完成，请确认是否发给 SuperBoss 评审",
     ["✓ 确认发送", "✗ 先不发"], True),
    (GREEN,  "任务完成",  "SuperCrew",
     "Milestone 4 所有接口开发完成，已推送 PR #23", [], False),
    (RED,    "异常上报",  "QA Agent",
     "验收发现：推荐结果中出现 2 个无效地点，已触发修复流程", [], False),
    (ACCENT, "进展摘要",  "EM Agent",
     "今日摘要：3 任务完成，1 待决策，无 blocker", [], False),
    (GRAY,   "自主决策",  "SuperBoss",
     "已自主决定使用方案 A（Redis 缓存），原因：性能更优", [], False),
]
for i, (clr, tag, actor, content, actions, urgent) in enumerate(feeds):
    y = 1.3 + i * 1.1
    bg_clr = RGBColor(0x22, 0x1a, 0x10) if urgent else CARD
    rect(slide, 1.7, y, 5.7, 0.95, bg_clr, radius=True)
    if urgent:
        rect(slide, 1.7, y, 0.06, 0.95, WARN)
    dot(slide, 1.9, y + 0.18, 0.14, clr)
    rect(slide, 2.12, y + 0.1, 1.0, 0.28, CARD_LT, radius=True)
    txt(slide, 2.15, y + 0.13, 0.94, 0.22, tag, size=10, color=clr,
        bold=True, align=PP_ALIGN.CENTER)
    txt(slide, 3.2, y + 0.1, 1.5, 0.25, actor, size=11, color=GRAY)
    txt(slide, 1.9, y + 0.46, 4.0, 0.38, content, size=12, color=LIGHT)
    for j, action in enumerate(actions):
        clr2 = ACCENT if j == 0 else CARD_LT
        rect(slide, 5.7 - len(actions) * 1.1 + j * 1.15, y + 0.55, 1.0, 0.32, clr2, radius=True)
        txt(slide, 5.72 - len(actions) * 1.1 + j * 1.15, y + 0.6, 0.96, 0.24,
            action, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ── Right: Decision Center ──
txt(slide, 7.7, 0.9, 5.3, 0.35, "决策中心", size=14, color=WHITE, bold=True)
txt(slide, 7.7, 1.28, 5.3, 0.28, "只有你能拍板的事", size=12, color=GRAY)
divider(slide, 7.7, 1.62, 5.3)

decisions = [
    (WARN, "产品方向", "MVP 要不要加离线模式？\nAgent 建议：不加，聚焦核心功能", "高"),
    (ACCENT, "对外发布", "周末去呀 v1.0 是否发布到 App Store？\n所有验收已通过", "中"),
]
for i, (clr, title, content, priority) in enumerate(decisions):
    y = 1.8 + i * 2.0
    rect(slide, 7.7, y, 5.3, 1.75, CARD, radius=True)
    rect(slide, 7.7, y, 0.06, 1.75, clr)
    rect(slide, 12.3, y + 0.15, 0.6, 0.3, CARD_LT, radius=True)
    txt(slide, 12.32, y + 0.18, 0.56, 0.24, priority, size=10, color=clr,
        bold=True, align=PP_ALIGN.CENTER)
    txt(slide, 7.95, y + 0.15, 4.2, 0.32, title, size=14, color=WHITE, bold=True)
    txt(slide, 7.95, y + 0.52, 4.2, 0.6, content, size=12, color=LIGHT)
    rect(slide, 7.95, y + 1.3, 1.1, 0.32, ACCENT, radius=True)
    txt(slide, 7.97, y + 1.34, 1.06, 0.26, "✓ 同意", size=11, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER)
    rect(slide, 9.15, y + 1.3, 1.1, 0.32, CARD_LT, radius=True)
    txt(slide, 9.17, y + 1.34, 1.06, 0.26, "↩ 修改", size=11, color=GRAY,
        bold=True, align=PP_ALIGN.CENTER)

# Interrupt control
rect(slide, 7.7, 5.9, 5.3, 1.3, RGBColor(0x1a, 0x10, 0x10), radius=True,
     border_color=RED)
txt(slide, 7.95, 6.0, 4.8, 0.3, "⚡  紧急中断", size=13, color=RED, bold=True)
txt(slide, 7.95, 6.35, 4.8, 0.45,
    "立即暂停指定 Agent 的当前任务\n中断指令走高优先级通道，实时生效",
    size=12, color=LIGHT)
rect(slide, 10.5, 6.08, 2.2, 0.88, RED, radius=True)
txt(slide, 10.52, 6.3, 2.16, 0.35, "暂停团队", size=13, color=WHITE,
    bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 5: Workflow — Natural Language → Agent Learns
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide)
sidebar(slide)
topbar(slide, "Workflows", "你不配置流程，你只需要告诉团队一次")
slide_label(slide, "Screen 5 of 5 — Workflow Canvas")

# ── Left: Natural Language Input ──
txt(slide, 1.85, 0.9, 5.5, 0.35, "告诉团队你的工作方式", size=14, color=WHITE, bold=True)

# Chat bubbles — user says something
rect(slide, 1.7, 1.35, 5.7, 3.8, CARD, radius=True)

messages = [
    ("🧑 Juanjuan", WARN, True,
     "以后 QA 验收必须附截图，\n没有截图不算通过"),
    ("🐙 pm-Octopus", ACCENT, False,
     "收到，已把验收需附截图写入\n自己的行为准则，QA Agent 同步更新。"),
    ("🧑 Juanjuan", WARN, True,
     "Bug issue 没有 repro steps\nDev 可以直接拒绝"),
    ("⚡ SuperCrew", PURPLE, False,
     "明白，已加入我的行为准则：\n无 repro 的 bug 我有权关闭。"),
]
for i, (name, clr, is_user, msg) in enumerate(messages):
    y = 1.5 + i * 0.82
    if is_user:
        rect(slide, 4.5, y, 2.7, 0.65, RGBColor(0x1a, 0x1e, 0x10), radius=True)
        txt(slide, 4.6, y + 0.06, 2.5, 0.2, name, size=10, color=clr, bold=True)
        txt(slide, 4.6, y + 0.28, 2.5, 0.35, msg, size=11, color=LIGHT)
    else:
        rect(slide, 1.85, y, 2.7, 0.65, CARD_LT, radius=True)
        txt(slide, 1.95, y + 0.06, 2.5, 0.2, name, size=10, color=clr, bold=True)
        txt(slide, 1.95, y + 0.28, 2.5, 0.35, msg, size=11, color=LIGHT)

# Input box
rect(slide, 1.7, 5.25, 5.7, 0.55, CARD_LT, radius=True, border_color=ACCENT)
txt(slide, 1.95, 5.35, 4.5, 0.35, "描述你的工作方式或流程规则...", size=12, color=DIM)
rect(slide, 6.8, 5.3, 0.5, 0.45, ACCENT, radius=True)
txt(slide, 6.82, 5.35, 0.46, 0.35, "↑", size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(slide, 1.7, 5.9, 5.7, 0.3,
    "Agent 自动理解、分配、写入各自行为准则",
    size=11, color=GREEN, align=PP_ALIGN.CENTER)

# ── Right: Agent Rules (Result) ──
txt(slide, 7.7, 0.9, 5.3, 0.35, "Agent 行为准则（自动更新）", size=14, color=WHITE, bold=True)
txt(slide, 7.7, 1.28, 5.3, 0.28, "这就是流程文档，始终最新", size=12, color=GRAY)

agent_rules = [
    ("🔍 QA Agent", WARN, [
        "验收必须附截图，截图需覆盖全部验收项",
        "验收不通过 → 打回 Dev，附截图说明原因",
        "重新验收需与原截图对比",
    ]),
    ("⚡ Dev Agent", PURPLE, [
        "接受 bug 前确认 repro steps，无法复现可关闭",
        "修复需附根因分析 + 单测",
        "收到无 repro 的 bug issue 有权拒绝",
    ]),
    ("🐙 PM Agent", ACCENT, [
        "需求变更先开 issue，不口头通知 Dev",
        "关 issue 前确认验收截图已存档",
    ]),
]
y_start = 1.65
for agent, clr, rules in agent_rules:
    rect(slide, 7.7, y_start, 5.3, 0.3, CARD_LT)
    txt(slide, 7.85, y_start + 0.05, 4, 0.22, agent, size=12, color=clr, bold=True)
    # updated badge
    rect(slide, 12.2, y_start + 0.05, 0.65, 0.22, RGBColor(0x10, 0x22, 0x18), radius=True)
    txt(slide, 12.22, y_start + 0.07, 0.61, 0.18, "已更新", size=9, color=GREEN, align=PP_ALIGN.CENTER)
    for j, rule in enumerate(rules):
        y = y_start + 0.35 + j * 0.38
        rect(slide, 7.7, y, 5.3, 0.33, CARD, radius=True)
        dot(slide, 7.85, y + 0.11, 0.1, clr)
        txt(slide, 8.05, y + 0.06, 4.8, 0.24, rule, size=11, color=LIGHT)
    y_start += 0.35 + len(rules) * 0.38 + 0.18

out = "/Users/juanjuanliu/.openclaw/workspace-pm-lead/tasks/hybrid-workforce-concept-2026-04-11/hybrid-workforce-ui.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
