#!/usr/bin/env python3
"""Hybrid Workforce — Pitch Deck"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Palette
BG       = RGBColor(0x0d, 0x0f, 0x1a)
CARD_BG  = RGBColor(0x16, 0x18, 0x2a)
CARD_LT  = RGBColor(0x1e, 0x21, 0x35)
ACCENT   = RGBColor(0x6c, 0x9e, 0xff)
ACCENT2  = RGBColor(0x7c, 0xdb, 0xc4)
WARN     = RGBColor(0xff, 0xb8, 0x6c)
GREEN    = RGBColor(0x6b, 0xe6, 0x9c)
WHITE    = RGBColor(0xf0, 0xf0, 0xf5)
GRAY     = RGBColor(0x8a, 0x8d, 0xa0)
DIM      = RGBColor(0x4a, 0x4d, 0x62)
LIGHT    = RGBColor(0xcc, 0xcc, 0xd5)
PURPLE   = RGBColor(0xb0, 0x8e, 0xff)
RED      = RGBColor(0xff, 0x7b, 0x7b)

def bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG

def txt(slide, l, t, w, h, text, size=18, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = "Helvetica Neue"; p.alignment = align
    return tf

def card(slide, l, t, w, h, fill=None):
    fill = fill or CARD_BG
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.adjustments[0] = 0.04
    return s

def line(slide, l, t, w, color=DIM):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(0.025))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

# ── SLIDE 1: Cover ────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide)

# Gradient-feel left bar
s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(0.5), Inches(7.5))
s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

txt(slide, 1.2, 1.5, 10, 0.6,
    "HYBRID WORKFORCE", size=14, color=ACCENT, bold=True)
txt(slide, 1.2, 2.1, 10, 1.6,
    "组织的下一个\n基本单位", size=52, color=WHITE, bold=True)
line(slide, 1.2, 4.2, 6, ACCENT)
txt(slide, 1.2, 4.5, 10, 0.6,
    "Build your workforce, not your tools.", size=20, color=LIGHT)
txt(slide, 1.2, 5.8, 10, 0.4,
    "2026", size=16, color=DIM)

# ── SLIDE 2: The Big Shift ────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide)

txt(slide, 1, 0.4, 11.3, 0.65,
    "一个正在发生的范式转移", size=32, color=WHITE, bold=True)

eras = [
    ("工业时代", "人 + 机器", "体力劳动外包给机器\n人负责操作和管理", GRAY),
    ("信息时代", "人 + 软件", "脑力劳动辅助给软件\n人负责决策和创造", GRAY),
    ("Agent 时代", "碳 + 硅团队", "执行本身外包给 Agent\n人负责方向和涌现", ACCENT),
]
for i, (era, unit, desc, clr) in enumerate(eras):
    x = 0.8 + i * 4.1
    is_now = i == 2
    card(slide, x, 1.4, 3.8, 4.2, CARD_LT if is_now else CARD_BG)
    if is_now:
        line(slide, x, 1.4, 3.8, ACCENT)
    txt(slide, x+0.25, 1.6, 3.3, 0.4, era, size=13, color=clr)
    txt(slide, x+0.25, 2.05, 3.3, 0.55, unit, size=24, color=WHITE if is_now else LIGHT, bold=is_now)
    line(slide, x+0.25, 2.7, 3.3, DIM)
    txt(slide, x+0.25, 2.9, 3.3, 1.5, desc, size=15, color=LIGHT)
    if is_now:
        card(slide, x+0.25, 4.4, 3.3, 0.7, DIM)
        txt(slide, x+0.35, 4.5, 3.1, 0.5,
            "✦  涌现 = 碳硅协作 > 双方之和", size=13, color=GREEN, bold=True)

card(slide, 0.8, 6.1, 11.7, 0.9, CARD_LT)
txt(slide, 1, 6.2, 11.3, 0.7,
    "未来的竞争力，不是\"会不会用 AI\"，而是\"能不能管理一支 AI 团队\"",
    size=18, color=WARN, bold=True, align=PP_ALIGN.CENTER)

# ── SLIDE 3: Why Now ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide)

txt(slide, 1, 0.4, 11.3, 0.65,
    "为什么是现在？四层技术同时成熟", size=32, color=WHITE, bold=True)

layers = [
    ("Harness Layer", "行为准则 · 权限边界 · 工作流约束",
     "Agent 的工作环境定义\n有了 harness，Agent 才有角色", ACCENT),
    ("Memory Layer", "决策记录 · 经验沉淀 · 上下文传承",
     "持久化的组织记忆\n知识不再只存在人脑里", ACCENT2),
    ("Learning Layer", "犯错 → 反思 → 规则更新 → 不再重犯",
     "Agent 从纠正中自我进化\n培训 Agent 就像培训员工", PURPLE),
    ("Orchestration Layer", "谁做什么 · 何时交接 · 卡住找谁",
     "碳硅之间的动态分工\n人和 Agent 协作有了协议", WARN),
]
for i, (title, subtitle, desc, clr) in enumerate(layers):
    x = 0.6 + i * 3.1
    card(slide, x, 1.4, 2.8, 4.5, CARD_BG)
    # Color top bar
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(1.4), Inches(2.8), Inches(0.08))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()
    txt(slide, x+0.2, 1.6, 2.4, 0.45, title, size=16, color=clr, bold=True)
    txt(slide, x+0.2, 2.1, 2.4, 0.45, subtitle, size=11, color=GRAY)
    line(slide, x+0.2, 2.6, 2.4, DIM)
    txt(slide, x+0.2, 2.8, 2.4, 1.8, desc, size=14, color=LIGHT)

txt(slide, 0.6, 6.1, 12, 0.4,
    "这四层同时成熟之前，AI 只能是工具。四层齐备，AI 才能成为同事。",
    size=17, color=GREEN, align=PP_ALIGN.CENTER)

# ── SLIDE 4: The Problem ─────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide)

txt(slide, 0.8, 0.25, 11, 1.6,
    "50 AI employees.\nZero management system.",
    size=40, color=WHITE, bold=True)

txt(slide, 0.8, 2.1, 2.2, 0.3, "THE PROBLEM", size=11, color=ACCENT, bold=True)
txt(slide, 6.5, 2.1, 6.5, 0.3, "WHAT ACTUALLY HAPPENED", size=11, color=GRAY, bold=True)

problems = [
    ("没有组织架构", "谁做什么、向谁汇报？"),
    ("没有组织记忆", "昨天的决策今天忘了"),
    ("没有培训体系", "同样的错反复犯，新人从零教起"),
    ("没有流程规范", "AI 自作主张，人类失控"),
]
for i, (title, sub) in enumerate(problems):
    y = 2.55 + i * 1.1
    txt(slide, 0.8, y, 5.5, 0.45, title, size=20, color=ACCENT, bold=True)
    txt(slide, 0.8, y + 0.45, 5.5, 0.35, sub, size=13, color=GRAY)

cases = [
    "PM 和 QA 都有验收职责，边界模糊，PM 看 QA 卡住就帮忙绕过",
    "PM 被冻结期间仍越界帮 QA 推文件、贴 issue",
    "Dev 没等需求评审就开工，PR 做完才发现方向错了",
    "连续 5 天无 daily notes，项目状态完全失明",
    "PR #101 验收停在 22%，无人知道是卡住还是被遗忘",
    "PM→QA 交接丢失决策上下文，实现违反 PRD 决策",
    "验收只看最近 PR，漏掉累积改动，通过后 bug 更多",
    "新 QA 加入当天：端口用错、流程不会、用代码审查代替实际运行",
    "新 PM 接手新项目，所有规则、流程都要再重复一遍",
    "AGENTS.md 膨胀到 20+ 条规则，关键规则被淹没，读了也记不住",
    "同一条规则纠正多次，跨 session 后又忘了，每个新 session 从零开始",
    "QA 没有截图就报验收 Pass，团队基于假结果做决策",
    "Agent 跑长任务时变成聋子，叫停指令等 2 小时才读到",
    "跑完错误任务后刷屏输出，淹没其他人的消息",
    "后端路径用错导致 LLM 走 fallback，整个上午验收结果不可靠",
]
for i, case in enumerate(cases):
    y = 2.22 + i * 0.32
    txt(slide, 6.5, y, 6.3, 0.28, f"· {case}", size=10.5, color=LIGHT)

# ── SLIDE 5_OLD: The Product ─────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide)

txt(slide, 1, 0.4, 11.3, 0.65,
    "Hybrid Workforce — 三个核心模块", size=32, color=WHITE, bold=True)

modules = [
    ("01", "Workforce\nBuilder", ACCENT,
     "搭建你的 AI 团队",
     "预设角色模板库\n开箱即用，30 分钟上线\n工具接入不绑定平台"),
    ("02", "Workforce\nTrainer", PURPLE,
     "培训和进化你的 Agent",
     "反馈 → 行为规则自动更新\n记忆体系跨天不失忆\n能力越用越强"),
    ("03", "Workforce\nOS", ACCENT2,
     "碳硅协作运作系统",
     "介入触发机制\nAgent 互相验收\n人只处理真正需要拍板的事"),
]
for i, (num, title, clr, sub, desc) in enumerate(modules):
    x = 0.8 + i * 4.1
    card(slide, x, 1.4, 3.8, 5.2, CARD_BG)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(1.4), Inches(3.8), Inches(0.08))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()
    txt(slide, x+0.25, 1.6, 0.6, 0.5, num, size=28, color=clr, bold=True)
    txt(slide, x+0.25, 2.2, 3.3, 0.8, title, size=22, color=WHITE, bold=True)
    txt(slide, x+0.25, 3.1, 3.3, 0.4, sub, size=14, color=WARN)
    line(slide, x+0.25, 3.55, 3.3, DIM)
    txt(slide, x+0.25, 3.75, 3.3, 2.4, desc, size=14, color=LIGHT)

# ── SLIDE 5: Proof — Real Case ────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide)

txt(slide, 1, 0.4, 11.3, 0.65,
    "我们自己就是第一个用户", size=32, color=WHITE, bold=True)
txt(slide, 1, 1.05, 11.3, 0.4,
    "周末去呀 — 一个真实的碳硅协作案例", size=18, color=GRAY)

# Left: human input
card(slide, 0.8, 1.7, 5.5, 1.1, CARD_BG)
txt(slide, 1.0, 1.8, 5.1, 0.35, "🧑  人的输入", size=13, color=WARN, bold=True)
txt(slide, 1.0, 2.15, 5.1, 0.5,
    "\"从小红书挖需求，做一个 MVP 产品\"  →  仅此而已",
    size=15, color=LIGHT)

# Results row
results = [
    ("需求挖掘", "4维评分\n选出周末去哪玩", ACCENT),
    ("PRD", "竞品分析+用户洞察\n完整需求文档", ACCENT2),
    ("开发", "5个Milestone\n前后端+部署", PURPLE),
    ("验收", "Agent互相验收\n发现LLM编数据", WARN),
]
for i, (title, desc, clr) in enumerate(results):
    x = 0.8 + i * 3.1
    card(slide, x, 3.1, 2.8, 1.9, CARD_BG)
    txt(slide, x+0.2, 3.2, 2.4, 0.35, title, size=15, color=clr, bold=True)
    line(slide, x+0.2, 3.6, 2.4, DIM)
    txt(slide, x+0.2, 3.75, 2.4, 1.1, desc, size=13, color=LIGHT)

# Key numbers
card(slide, 0.8, 5.3, 11.7, 1.5, CARD_LT)
txt(slide, 1.0, 5.45, 11.3, 0.35, "关键数字", size=14, color=GRAY)
nums = [("1 句话", "人的全部输入"), ("5 次换人", "项目零中断"),
        ("92 个真实 POI", "Agent 自主建立防线后"), ("0 次全程跟进", "人不需要实时在线")]
for i, (n, l) in enumerate(nums):
    x = 1.0 + i * 3.0
    txt(slide, x, 5.85, 2.7, 0.45, n, size=20, color=GREEN, bold=True)
    txt(slide, x, 6.3, 2.7, 0.35, l, size=12, color=GRAY)

# ── SLIDE 6: New Roles ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide)

txt(slide, 1, 0.4, 11.3, 0.65,
    "新职业物种正在出现", size=32, color=WHITE, bold=True)
txt(slide, 1, 1.05, 8, 0.4,
    "这些角色现在没有名字，但已经存在了", size=18, color=GRAY)

roles = [
    ("Harness\nEngineer", ACCENT,
     "设计 Agent 的行为边界、角色准则、工作流约束\n让 Agent 成为可靠的团队成员，而不是失控的工具"),
    ("AI\nTeam Lead", PURPLE,
     "管理碳硅混编团队\n懂得何时授权给 Agent、何时人必须介入"),
    ("Context\nArchitect", ACCENT2,
     "构建组织的记忆和知识体系\n让知识不依赖任何单一个体，换人不断档"),
]
for i, (title, clr, desc) in enumerate(roles):
    x = 0.8 + i * 4.1
    card(slide, x, 1.7, 3.8, 4.0, CARD_BG)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(1.7), Inches(3.8), Inches(0.08))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()
    txt(slide, x+0.25, 1.9, 3.3, 0.8, title, size=26, color=clr, bold=True)
    line(slide, x+0.25, 2.8, 3.3, DIM)
    txt(slide, x+0.25, 3.0, 3.3, 2.4, desc, size=14, color=LIGHT)

card(slide, 0.8, 6.0, 11.7, 0.95, CARD_LT)
txt(slide, 1.0, 6.1, 11.3, 0.75,
    "我们在做的，不是一个 AI 工具产品。\n我们在定义一个新的工作范式，以及支撑这个范式的基础设施。",
    size=16, color=WARN, align=PP_ALIGN.CENTER)

# ── SLIDE 7: The Trust Layer ─────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide)

txt(slide, 1, 0.3, 11.3, 0.65,
    "别人的 AI 说完就忘，我们的 Agent 是真正在成长",
    size=30, color=WHITE, bold=True)
txt(slide, 1, 0.95, 11.3, 0.35,
    "核心差异化：学到的东西可信 · 可见 · 可追溯",
    size=16, color=ACCENT)

# 2x2 pain point cards with multiple real cases each
pains = [
    (RED, "同样的错，反复犯",
     [
         "LLM 编数据通过了验收，上线后才发现全是假 POI",
         "纠正后下次任务照犯，因为纠正没有真正写进规则",
         "PM 和 QA 都有验收职责，边界模糊，互相踢皮球",
         "PM 看 QA 卡住直接绕过，验收流程形同虚设",
     ]),
    (WARN, "说完就忘，无法验证",
     [
         "Agent 回复「已记住」，关掉对话重开是全新 Agent",
         "没有任何机制验证规则是否真正落盘",
         "人以为 Agent 学会了，实际只在当次对话里有效",
         "叫停指令因消息队列延迟，Agent 收到时已完成任务",
     ]),
    (PURPLE, "新 Agent 入职，重新教一遍",
     [
         "Dev Agent 换了一个，两周磨合的规则全部归零",
         "没有标准化 onboarding，每次靠人口头传递经验",
         "新 Agent 不知道历史决策，重复走弯路",
         "团队规模扩大时，培训成本线性增长",
     ]),
    (GRAY, "进展不透明，决策无据可查",
     [
         "不知道 Agent 在干什么，不放心又不想一直盯",
         "Agent 做了决定，不知道依据什么、有无参考规则",
         "多 Agent 并行时，人无法追溯协作过程",
         "出问题时无法定位：是规则问题还是执行问题",
     ]),
]
positions = [(0.8, 1.55), (6.82, 1.55), (0.8, 4.15), (6.82, 4.15)]
for (x, y), (clr, title, cases) in zip(positions, pains):
    card(slide, x, y, 5.72, 2.3, CARD_BG)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(5.72), Inches(0.055))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()
    txt(slide, x+0.2, y+0.12, 5.3, 0.35, title, size=16, color=clr, bold=True)
    for j, case in enumerate(cases):
        txt(slide, x+0.2, y+0.55+j*0.4, 5.3, 0.35, f"· {case}", size=11, color=LIGHT)

# Bottom 3 pillars
for i, (title, clr, desc) in enumerate([
    ("可信", GREEN,   "规则强制落盘，系统验证，不靠 Agent 自觉"),
    ("可见", ACCENT,  "Agent Profile 展示所有已学规则 + 学习时间"),
    ("可追溯", PURPLE, "执行时标注依据规则，决策有据可查"),
]):
    x = 0.8 + i * 4.18
    card(slide, x, 6.65, 3.9, 0.72, CARD_LT)
    txt(slide, x+0.2, 6.73, 1.1, 0.38, title, size=18, color=clr, bold=True)
    txt(slide, x+1.4, 6.79, 2.6, 0.3, desc, size=11, color=GRAY)

# ── SLIDE 8: Vision ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide)

s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(0.5), Inches(7.5))
s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

txt(slide, 1.2, 1.4, 10, 0.5, "HYBRID WORKFORCE", size=14, color=ACCENT, bold=True)
txt(slide, 1.2, 2.0, 10, 2.2,
    "过去 100 年，\n组织的基本单位是\"人\"。\n未来 10 年，是\"碳硅团队\"。",
    size=34, color=WHITE, bold=True)
line(slide, 1.2, 4.6, 7, ACCENT)
txt(slide, 1.2, 4.9, 10, 0.5,
    "Build your workforce, not your tools.",
    size=22, color=LIGHT)
txt(slide, 1.2, 6.1, 10, 0.4,
    "Juanjuan & Team  ·  2026", size=15, color=DIM)

out = "/Users/juanjuanliu/.openclaw/workspace-pm-lead/tasks/hybrid-workforce-concept-2026-04-11/hybrid-workforce-pitch.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
