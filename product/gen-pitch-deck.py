#!/usr/bin/env python3
"""Generate Hybrid Workforce Pitch Deck (.pptx)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# === Config ===
BG = RGBColor(10, 10, 20)
WHITE = RGBColor(255, 255, 255)
BODY = RGBColor(200, 200, 210)
PURPLE = RGBColor(108, 92, 231)
BLUE = RGBColor(0, 210, 255)
FONT = "Calibri"

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def add_slide():
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide


def add_text(slide, text, left, top, width, height, font_size=18, color=BODY,
             bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_title(slide, text, top=Inches(0.8)):
    add_text(slide, text, Inches(1), top, Inches(11.3), Inches(1),
             font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)


def add_bullets(slide, items, top=Inches(2.5), left=Inches(1.2), color=BODY,
                font_size=20, spacing=Inches(0.6), accent_prefix=True):
    for i, item in enumerate(items):
        y = top + Inches(i * 0.7)
        if accent_prefix and "——" in item:
            parts = item.split("——", 1)
            txBox = slide.shapes.add_textbox(left, y, Inches(10.5), Inches(0.6))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run1 = p.add_run()
            run1.text = parts[0] + " "
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = PURPLE
            run1.font.bold = True
            run1.font.name = FONT
            run2 = p.add_run()
            run2.text = "— " + parts[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = FONT
        else:
            add_text(slide, "▸  " + item, left, y, Inches(10.5), Inches(0.6),
                     font_size=font_size, color=color)


def add_purple_bar(slide, top, left, width, height):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = PURPLE
    shape.line.fill.background()
    return shape


# ============================================================
# Slide 1: Cover
# ============================================================
s = add_slide()
add_text(s, "HYBRID WORKFORCE", Inches(1), Inches(2.2), Inches(11.3), Inches(1.2),
         font_size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, "The Operating System for Human-AI Teams", Inches(1), Inches(3.5), Inches(11.3), Inches(0.8),
         font_size=28, color=PURPLE, bold=False, alignment=PP_ALIGN.CENTER)
# decorative line
add_purple_bar(s, Inches(4.5), Inches(5.5), Inches(2.3), Inches(0.04))
add_text(s, "碳硅混合团队的操作系统", Inches(1), Inches(5.0), Inches(11.3), Inches(0.6),
         font_size=18, color=BODY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 2: The Shift
# ============================================================
s = add_slide()
add_title(s, "AI is no longer a tool.\nIt's your coworker.")
add_bullets(s, [
    "Agent 从聊天助手进化为自主执行者",
    "Harness Engineering 成为新学科——管理 AI 比训练 AI 更重要",
    "每个知识工作者都将拥有自己的 AI 团队",
], top=Inches(3.0), accent_prefix=False)

# ============================================================
# Slide 3: The Problem (two-column: problem + real cases)
# ============================================================
s = add_slide()
add_title(s, "50 AI employees.\nZero management system.")

# Column headers
add_text(s, "THE PROBLEM", Inches(1), Inches(2.7), Inches(4.5), Inches(0.3),
         font_size=11, color=PURPLE, bold=True)
add_text(s, "WHAT ACTUALLY HAPPENED", Inches(6), Inches(2.7), Inches(6.5), Inches(0.3),
         font_size=11, color=RGBColor(120, 120, 140), bold=True)

# Row data: (problem, subtitle, [case1, case2, ...])
rows = [
    ("没有组织架构", "谁做什么、向谁汇报？", [
        "PM 和 QA 都有验收职责，边界模糊，PM 看 QA 卡住就帮忙绕过",
        "PM 被冻结期间仍越界帮 QA 推文件、贴 issue",
        "Dev 没等需求评审就开工，PR 做完才发现方向错了",
    ]),
    ("没有组织记忆", "昨天的决策今天忘了", [
        "连续 5 天无 daily notes，项目状态完全失明",
        "PR #101 验收停在 22%，无人知道是卡住还是被遗忘",
        "PM→QA 交接丢失决策上下文，实现违反 PRD 决策",
        "验收只看最近 PR，漏掉累积改动，通过后 bug 更多",
    ]),
    ("存储 ≠ 召回", "Teach Once, Store Forever, Recall Never", [
        "Manta 的 Skill 里写了「GitHub issue 图片用 markdown 格式」，但每次都忘，需要反复提醒",
        "规则不是没有，是在 20+ 条中被上下文稀释——Agent 不会像人一样「触景生情」",
        "同一条规则纠正多次，跨 session 后又忘了，每个新 session 从零开始",
        "新 QA 加入当天：端口用错、流程不会、用代码审查代替实际运行",
    ]),
    ("没有流程规范", "AI 自作主张，人类失控", [
        "QA 没有截图就报验收 Pass，团队基于假结果做决策",
        "Agent 跑长任务时变成聋子，叫停指令等 2 小时才读到",
        "跑完错误任务后刷屏输出，淹没其他人的消息",
        "后端路径用错导致 LLM 走 fallback，整个上午验收结果不可靠",
    ]),
]

CASE_COLOR = RGBColor(130, 130, 150)

for i, (prob, sub, cases) in enumerate(rows):
    y_base = Inches(3.05 + i * 1.05)

    # Left: problem (purple bold + gray sub)
    txBox = s.shapes.add_textbox(Inches(1), y_base, Inches(4.5), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = prob
    run1.font.size = Pt(18)
    run1.font.color.rgb = PURPLE
    run1.font.bold = True
    run1.font.name = FONT
    p2 = tf.add_paragraph()
    p2.text = sub
    p2.font.size = Pt(13)
    p2.font.color.rgb = BODY
    p2.font.name = FONT

    # Right: multiple real cases (small, compact)
    case_text = "\n".join(["· " + c for c in cases])
    add_text(s, case_text, Inches(6), y_base, Inches(6.5), Inches(0.95),
             font_size=11, color=CASE_COLOR)

    # Subtle divider
    if i < len(rows) - 1:
        div_y = Inches(3.05 + (i + 1) * 1.05 - 0.08)
        add_purple_bar(s, div_y, Inches(1), Inches(11.5), Inches(0.004))

# ============================================================
# Slide 4: The Vision
# ============================================================
s = add_slide()
add_text(s, "HYBRID WORKFORCE", Inches(1), Inches(1.5), Inches(11.3), Inches(1),
         font_size=24, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, "The Operating System\nfor Human-AI Teams", Inches(1), Inches(2.5), Inches(11.3), Inches(1.5),
         font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_purple_bar(s, Inches(4.3), Inches(5.5), Inches(2.3), Inches(0.04))
add_text(s, "让每个知识工作者能像管理真人团队一样，\n组建、培训、管理一支 AI 团队。",
         Inches(1.5), Inches(4.8), Inches(10.3), Inches(1.2),
         font_size=22, color=BODY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 5: Four Layers
# ============================================================
s = add_slide()
add_title(s, "Four Layers of Hybrid Workforce", top=Inches(0.6))

layers = [
    ("Harness Layer", "管理制度", "行为准则、权限边界"),
    ("Memory Layer", "组织记忆", "决策记录、经验沉淀"),
    ("Learning Layer", "员工成长", "纠正→规则→进化"),
    ("Orchestration Layer", "项目管理", "碳硅分工、流程调度"),
]
box_w = Inches(2.6)
box_h = Inches(3.2)
gap = Inches(0.35)
start_x = Inches(0.9)
top_y = Inches(2.3)

for i, (name, cn, desc) in enumerate(layers):
    x = start_x + (box_w + gap) * i
    shape = s.shapes.add_shape(1, x, top_y, box_w, box_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(20, 20, 35)
    shape.line.color.rgb = PURPLE
    shape.line.width = Pt(1.5)
    # layer name
    add_text(s, name, x + Inches(0.2), top_y + Inches(0.4), box_w - Inches(0.4), Inches(0.5),
             font_size=18, color=PURPLE, bold=True)
    add_text(s, cn, x + Inches(0.2), top_y + Inches(1.0), box_w - Inches(0.4), Inches(0.4),
             font_size=22, color=WHITE, bold=True)
    add_text(s, desc, x + Inches(0.2), top_y + Inches(1.7), box_w - Inches(0.4), Inches(1.2),
             font_size=16, color=BODY)

# ============================================================
# Slide 6: Build
# ============================================================
s = add_slide()
add_text(s, "PHASE 1", Inches(1), Inches(0.5), Inches(2), Inches(0.4),
         font_size=14, color=PURPLE, bold=True)
add_title(s, "Build Your AI Team", top=Inches(0.9))
add_bullets(s, [
    "团队模板——一键拉起 PM + Dev + QA",
    "角色定义——每个 Agent 有身份、准则、技能",
    "共享知识库——团队共享上下文，不重复交代",
], top=Inches(2.5))
add_text(s, "30 分钟从零到开始干活", Inches(1.2), Inches(5.2), Inches(10), Inches(0.5),
         font_size=24, color=BLUE, bold=True)

# ============================================================
# Slide 7: Train
# ============================================================
s = add_slide()
add_text(s, "PHASE 2", Inches(1), Inches(0.5), Inches(2), Inches(0.4),
         font_size=14, color=PURPLE, bold=True)
add_title(s, "Teach Once, Recall Instantly, Remember Forever", top=Inches(0.9))
add_text(s, "Today: Teach Once, Store Forever, Recall Never.\nTomorrow: Teach Once, Recall Instantly, Remember Forever.",
         Inches(1), Inches(1.6), Inches(10), Inches(0.7),
         font_size=18, color=BLUE, bold=False)
add_bullets(s, [
    "场景触发式召回——Agent 执行工具前，系统自动注入精准规则（pre-tool-call hook）",
    "你说一句话 → 每个 Agent 自己识别相关规则 → 写入行为准则",
    "四层记忆保障：场景触发 / 行为前检查 / 失败记忆召回 / 渐进式强化",
    "不改模型，改 Harness——轻量、即时、可解释",
], top=Inches(2.8))

# ============================================================
# Slide 8: Flow
# ============================================================
s = add_slide()
add_text(s, "PHASE 3", Inches(1), Inches(0.5), Inches(2), Inches(0.4),
         font_size=14, color=PURPLE, bold=True)
add_title(s, "Orchestrate Human-AI Workflows", top=Inches(0.9))
add_bullets(s, [
    "自然语言定义流程——说一遍，Agent 各自领会并执行",
    "Human Gate——关键决策必须人类拍板",
    "Kill Switch——实时中断，人类永远掌控",
    "Live Monitor——随时知道 AI 在做什么",
], top=Inches(2.5))

# ============================================================
# Slide 9: Why Us
# ============================================================
s = add_slide()
add_title(s, "We Are User Zero")
add_bullets(s, [
    "我们自己就是第一批用户",
    "一个人类 PM 指挥一支 AI 团队，从 0 到 MVP",
    "这不是理论，是我们每天在做的事",
    "所有痛点都来自真实实践",
], top=Inches(2.5), accent_prefix=False)

# ============================================================
# Slide 10: What's Next
# ============================================================
s = add_slide()
add_text(s, "The Future of Work is Hybrid", Inches(1), Inches(1.5), Inches(11.3), Inches(0.8),
         font_size=24, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s, "Are you ready to\nbuild your AI team?", Inches(1), Inches(2.8), Inches(11.3), Inches(1.8),
         font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_purple_bar(s, Inches(5.0), Inches(5.5), Inches(2.3), Inches(0.04))
add_text(s, "contact@hybridworkforce.ai  ·  hybridworkforce.ai",
         Inches(1), Inches(5.8), Inches(11.3), Inches(0.5),
         font_size=16, color=BODY, alignment=PP_ALIGN.CENTER)

# === Save ===
out = os.path.join(os.path.dirname(__file__), "hybrid-workforce-pitch.pptx")
prs.save(out)
print(f"✅ Saved to {out}")
